import os

import cv2
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tqdm import tqdm
from tempfile import TemporaryDirectory
from paddleocr import img_decode


class PptxParser:

    def __init__(self, ocr):
        self.ocr = ocr

    def __save_image__(self, block, save_folder):
        img_path = os.path.join(
            save_folder, "{}.jpg".format(block.sha1)
        )
        cv2.imwrite(img_path, img_decode(block.blob))
        return img_path

    def __extract(self, shape, save_folder):
        if shape.shape_type == 19:
            table = shape.table
            rows = ["<table><thead>"]
            row_text = ["<tr>"]
            for cell in table.rows[0].cells:
                row_text.append(f"<td>{cell.text_frame.text}</td>")
            row_text.append("</tr>")
            rows.append("".join(row_text))
            rows.append("</thead><tbody>")
            for idx, row in enumerate(table.rows):
                if idx == 0:
                    continue
                row_text = ["<tr>"]
                for cell in row.cells:
                    row_text.append(f"<td>{cell.text_frame.text}</td>")
                row_text.append("</tr>")
                rows.append("".join(row_text))
            rows.append("</tbody></table>")
            s = "以下是一张表格，包含对应的格式以及数据，以<end_table>结尾: \n"
            table = "".join(rows)
            s += table
            s += "\n<end_table>"
            return s

        if shape.has_text_frame:
            return shape.text_frame.text

        if shape.shape_type == 6:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                t = self.__extract(p, save_folder)
                if t:
                    texts.append(t)
            return "\n".join(texts)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            res = self.ocr(shape.image.blob)
            if len(res):
                return "\n".join(res)

    def __call__(self, pptx_path):
        prs = Presentation(pptx_path)
        re = []
        with TemporaryDirectory() as tmpdir:
            for idx, slide in enumerate(tqdm(prs.slides, desc=f"process {pptx_path}")):
                texts = []
                for shape in tqdm(sorted(slide.shapes, key=lambda x: (x.top // 10, x.left)), desc=f"page {idx}"):
                    txt = self.__extract(shape, tmpdir)
                    if txt:
                        texts.append(txt)
                re.append("\n".join(texts))
        return re


if __name__ == '__main__':
    from pdf_parser import PdfParser
    re = PptxParser(PdfParser())('../data/pptx/LargePpt.pptx')
    with open('../output/LargePpt.txt', 'w') as w:
        w.write("\n".join(re))
